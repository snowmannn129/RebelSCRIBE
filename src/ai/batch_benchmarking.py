#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Batch Benchmarking for RebelSCRIBE.

This module provides functionality for running multiple benchmarks in batch mode,
managing benchmark templates, and generating comprehensive reports.

Example usage:
    ```python
    from src.ai.batch_benchmarking import (
        BatchBenchmark, BenchmarkTemplate, run_batch_benchmark,
        create_benchmark_template, get_benchmark_templates,
        export_batch_results
    )
    
    # Create a benchmark template
    template = create_benchmark_template(
        name="Creative Writing",
        prompts=[
            "Write a short story about a robot learning to paint.",
            "Write a poem about the changing seasons.",
            "Write a dialogue between two characters discussing the ethics of AI."
        ],
        max_tokens=150,
        num_runs=3,
        tags=["creative", "writing"]
    )
    
    # Create a batch benchmark
    batch = BatchBenchmark(
        name="Creative Writing Benchmark",
        model_ids=["llama-2-7b", "mistral-7b", "phi-2"],
        template_id=template.id,
        description="Comparing models on creative writing tasks"
    )
    
    # Run the batch benchmark
    results = run_batch_benchmark(batch)
    
    # Export the results
    export_batch_results(results, format="pdf", output_path="benchmark_results.pdf")
    ```
"""

import os
import json
import time
import logging
import threading
import statistics
import uuid
import tempfile
from enum import Enum, auto
from typing import Dict, List, Optional, Any, Set, Tuple, Union, Callable
from dataclasses import dataclass, field, asdict
import datetime
import concurrent.futures
from pathlib import Path

from src.utils.logging_utils import get_logger
from src.utils.file_utils import ensure_directory
from src.ai.progress_callbacks import (
    OperationType, ProgressStatus, ProgressInfo, ProgressCallback,
    create_operation, start_operation, update_progress,
    complete_operation, fail_operation
)
from src.ai.model_benchmarking import (
    ModelBenchmark, BenchmarkResult, run_benchmark,
    get_benchmark_results, BenchmarkRegistry
)
from src.ai.model_registry import (
    ModelRegistry, ModelInfo, ModelType, ModelSource
)
from src.ai.benchmark_visualization import (
    plot_model_comparison, plot_benchmark_history, plot_metric_correlation,
    plot_model_radar_chart, export_visualization, generate_benchmark_report_html,
    check_visualization_dependencies
)

logger = get_logger(__name__)


class BatchStatus(Enum):
    """Enum representing the status of a batch benchmark."""
    PENDING = auto()
    RUNNING = auto()
    COMPLETED = auto()
    FAILED = auto()
    CANCELLED = auto()


@dataclass
class BenchmarkTemplate:
    """Class for defining a benchmark template."""
    id: str
    name: str
    prompts: List[str]
    max_tokens: int = 100
    num_runs: int = 3
    temperature: float = 0.7
    top_p: float = 0.9
    tags: List[str] = field(default_factory=list)
    description: Optional[str] = None
    reference_texts: Optional[List[str]] = None
    save_logprobs: bool = True
    created_at: str = field(default_factory=lambda: datetime.datetime.now().isoformat())
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the template to a dictionary."""
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BenchmarkTemplate':
        """Create a BenchmarkTemplate instance from a dictionary."""
        return cls(**data)


@dataclass
class BatchBenchmark:
    """Class for defining a batch benchmark."""
    id: str
    name: str
    model_ids: List[str]
    template_id: str
    description: Optional[str] = None
    tags: List[str] = field(default_factory=list)
    parallel: bool = False
    max_workers: int = 1
    status: BatchStatus = BatchStatus.PENDING
    created_at: str = field(default_factory=lambda: datetime.datetime.now().isoformat())
    started_at: Optional[str] = None
    completed_at: Optional[str] = None
    error: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the batch benchmark to a dictionary."""
        data = asdict(self)
        # Convert enum to string
        data["status"] = self.status.name
        return data
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BatchBenchmark':
        """Create a BatchBenchmark instance from a dictionary."""
        # Convert string status back to enum
        status_str = data.pop("status", "PENDING")
        status = BatchStatus[status_str]
        
        return cls(
            status=status,
            **data
        )


@dataclass
class BatchResult:
    """Class for storing batch benchmark results."""
    batch_id: str
    name: str
    template_id: str
    model_ids: List[str]
    benchmark_results: Dict[str, List[BenchmarkResult]]
    started_at: str
    completed_at: str
    duration_seconds: float
    error: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the batch result to a dictionary."""
        data = asdict(self)
        # Convert BenchmarkResult objects to dictionaries
        data["benchmark_results"] = {
            model_id: [result.to_dict() for result in results]
            for model_id, results in self.benchmark_results.items()
        }
        return data
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BatchResult':
        """Create a BatchResult instance from a dictionary."""
        # Convert dictionaries back to BenchmarkResult objects
        benchmark_results_dict = data.pop("benchmark_results", {})
        benchmark_results = {
            model_id: [BenchmarkResult.from_dict(result_dict) for result_dict in results_dict]
            for model_id, results_dict in benchmark_results_dict.items()
        }
        
        return cls(
            benchmark_results=benchmark_results,
            **data
        )


class BatchBenchmarkRegistry:
    """
    Singleton class for managing batch benchmarks, templates, and results.
    
    This class provides functionality for storing, retrieving, and managing
    batch benchmark definitions, templates, and results.
    """
    
    _instance = None
    _lock = threading.Lock()
    
    @classmethod
    def get_instance(cls) -> 'BatchBenchmarkRegistry':
        """Get the singleton instance of the BatchBenchmarkRegistry."""
        with cls._lock:
            if cls._instance is None:
                cls._instance = cls()
            return cls._instance
    
    def __init__(self):
        """Initialize the batch benchmark registry."""
        # Ensure this is only called once
        if BatchBenchmarkRegistry._instance is not None:
            raise RuntimeError("Use BatchBenchmarkRegistry.get_instance() to get the singleton instance")
        
        # Initialize registry data
        self._templates: Dict[str, BenchmarkTemplate] = {}
        self._batches: Dict[str, BatchBenchmark] = {}
        self._results: Dict[str, BatchResult] = {}
        self._registry_file = self._get_registry_file_path()
        
        # Load existing registry data
        self._load_registry()
    
    def _get_registry_file_path(self) -> str:
        """Get the path to the registry file."""
        # Get the registry directory
        registry_dir = os.environ.get(
            "REBELSCRIBE_BENCHMARK_DIR", 
            os.path.join(os.path.expanduser("~"), ".rebelscribe", "benchmarks")
        )
        
        # Ensure the directory exists
        ensure_directory(registry_dir)
        
        # Return the path to the registry file
        return os.path.join(registry_dir, "batch_benchmark_registry.json")
    
    def _load_registry(self) -> None:
        """Load the registry data from the registry file."""
        if not os.path.exists(self._registry_file):
            logger.info(f"Registry file not found at {self._registry_file}, creating new registry")
            return
        
        try:
            with open(self._registry_file, "r") as f:
                data = json.load(f)
            
            # Convert dictionaries to BenchmarkTemplate objects
            for template_id, template_data in data.get("templates", {}).items():
                self._templates[template_id] = BenchmarkTemplate.from_dict(template_data)
            
            # Convert dictionaries to BatchBenchmark objects
            for batch_id, batch_data in data.get("batches", {}).items():
                self._batches[batch_id] = BatchBenchmark.from_dict(batch_data)
            
            # Convert dictionaries to BatchResult objects
            for result_id, result_data in data.get("results", {}).items():
                self._results[result_id] = BatchResult.from_dict(result_data)
            
            logger.info(
                f"Loaded {len(self._templates)} templates, {len(self._batches)} batches, "
                f"and {len(self._results)} results from registry"
            )
        
        except Exception as e:
            logger.error(f"Error loading registry: {e}", exc_info=True)
            # Create a backup of the corrupted file
            if os.path.exists(self._registry_file):
                backup_path = f"{self._registry_file}.bak.{int(time.time())}"
                try:
                    import shutil
                    shutil.copy2(self._registry_file, backup_path)
                    logger.info(f"Created backup of corrupted registry at {backup_path}")
                except Exception as backup_error:
                    logger.error(f"Error creating backup: {backup_error}")
    
    def _save_registry(self) -> None:
        """Save the registry data to the registry file."""
        try:
            # Convert objects to dictionaries
            templates_dict = {template_id: template.to_dict() for template_id, template in self._templates.items()}
            batches_dict = {batch_id: batch.to_dict() for batch_id, batch in self._batches.items()}
            results_dict = {result_id: result.to_dict() for result_id, result in self._results.items()}
            
            # Create the data to save
            data = {
                "templates": templates_dict,
                "batches": batches_dict,
                "results": results_dict,
                "last_updated": datetime.datetime.now().isoformat()
            }
            
            # Save to file
            with open(self._registry_file, "w") as f:
                json.dump(data, f, indent=2)
            
            logger.info(
                f"Saved {len(self._templates)} templates, {len(self._batches)} batches, "
                f"and {len(self._results)} results to registry"
            )
        
        except Exception as e:
            logger.error(f"Error saving registry: {e}", exc_info=True)
    
    def register_template(self, template: BenchmarkTemplate) -> str:
        """
        Register a benchmark template in the registry.
        
        Args:
            template: The benchmark template to register.
            
        Returns:
            str: The ID of the registered template.
        """
        with self._lock:
            # Use the provided ID or generate a new one
            template_id = template.id
            
            # Register the template
            self._templates[template_id] = template
            self._save_registry()
            
            logger.info(f"Registered benchmark template: {template_id}")
            return template_id
    
    def get_template(self, template_id: str) -> Optional[BenchmarkTemplate]:
        """
        Get a benchmark template by ID.
        
        Args:
            template_id: The ID of the template to get.
            
        Returns:
            Optional[BenchmarkTemplate]: The template, or None if not found.
        """
        return self._templates.get(template_id)
    
    def get_all_templates(self) -> List[BenchmarkTemplate]:
        """
        Get all registered benchmark templates.
        
        Returns:
            List[BenchmarkTemplate]: A list of benchmark templates.
        """
        return list(self._templates.values())
    
    def register_batch(self, batch: BatchBenchmark) -> str:
        """
        Register a batch benchmark in the registry.
        
        Args:
            batch: The batch benchmark to register.
            
        Returns:
            str: The ID of the registered batch.
        """
        with self._lock:
            # Use the provided ID or generate a new one
            batch_id = batch.id
            
            # Register the batch
            self._batches[batch_id] = batch
            self._save_registry()
            
            logger.info(f"Registered batch benchmark: {batch_id}")
            return batch_id
    
    def update_batch(self, batch: BatchBenchmark) -> None:
        """
        Update a batch benchmark in the registry.
        
        Args:
            batch: The batch benchmark to update.
        """
        with self._lock:
            # Update the batch
            self._batches[batch.id] = batch
            self._save_registry()
            
            logger.info(f"Updated batch benchmark: {batch.id}")
    
    def get_batch(self, batch_id: str) -> Optional[BatchBenchmark]:
        """
        Get a batch benchmark by ID.
        
        Args:
            batch_id: The ID of the batch to get.
            
        Returns:
            Optional[BatchBenchmark]: The batch, or None if not found.
        """
        return self._batches.get(batch_id)
    
    def get_all_batches(self) -> List[BatchBenchmark]:
        """
        Get all registered batch benchmarks.
        
        Returns:
            List[BatchBenchmark]: A list of batch benchmarks.
        """
        return list(self._batches.values())
    
    def register_result(self, result: BatchResult) -> str:
        """
        Register a batch result in the registry.
        
        Args:
            result: The batch result to register.
            
        Returns:
            str: The ID of the registered result.
        """
        with self._lock:
            # Use the batch ID as the result ID
            result_id = result.batch_id
            
            # Register the result
            self._results[result_id] = result
            self._save_registry()
            
            logger.info(f"Registered batch result: {result_id}")
            return result_id
    
    def get_result(self, result_id: str) -> Optional[BatchResult]:
        """
        Get a batch result by ID.
        
        Args:
            result_id: The ID of the result to get.
            
        Returns:
            Optional[BatchResult]: The result, or None if not found.
        """
        return self._results.get(result_id)
    
    def get_all_results(self) -> List[BatchResult]:
        """
        Get all registered batch results.
        
        Returns:
            List[BatchResult]: A list of batch results.
        """
        return list(self._results.values())


def create_benchmark_template(
    name: str,
    prompts: List[str],
    max_tokens: int = 100,
    num_runs: int = 3,
    temperature: float = 0.7,
    top_p: float = 0.9,
    tags: Optional[List[str]] = None,
    description: Optional[str] = None,
    reference_texts: Optional[List[str]] = None,
    save_logprobs: bool = True
) -> BenchmarkTemplate:
    """
    Create a benchmark template.
    
    Args:
        name: The name of the template.
        prompts: The list of prompts to use for benchmarking.
        max_tokens: The maximum number of tokens to generate.
        num_runs: The number of times to run each benchmark.
        temperature: The temperature to use for generation.
        top_p: The top-p value to use for generation.
        tags: Optional list of tags for the template.
        description: Optional description of the template.
        reference_texts: Optional list of reference texts for BLEU score calculation.
        save_logprobs: Whether to save token logprobs for perplexity calculation.
        
    Returns:
        BenchmarkTemplate: The created benchmark template.
    """
    # Generate a unique ID for the template
    template_id = f"template_{uuid.uuid4().hex[:8]}"
    
    # Create the template
    template = BenchmarkTemplate(
        id=template_id,
        name=name,
        prompts=prompts,
        max_tokens=max_tokens,
        num_runs=num_runs,
        temperature=temperature,
        top_p=top_p,
        tags=tags or [],
        description=description,
        reference_texts=reference_texts,
        save_logprobs=save_logprobs
    )
    
    # Register the template
    registry = BatchBenchmarkRegistry.get_instance()
    registry.register_template(template)
    
    logger.info(f"Created benchmark template: {template_id}")
    return template


def get_benchmark_templates() -> List[BenchmarkTemplate]:
    """
    Get all registered benchmark templates.
    
    Returns:
        List[BenchmarkTemplate]: A list of benchmark templates.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_templates()


def create_batch_benchmark(
    name: str,
    model_ids: List[str],
    template_id: str,
    description: Optional[str] = None,
    tags: Optional[List[str]] = None,
    parallel: bool = False,
    max_workers: int = 1
) -> BatchBenchmark:
    """
    Create a batch benchmark.
    
    Args:
        name: The name of the batch benchmark.
        model_ids: The list of model IDs to benchmark.
        template_id: The ID of the template to use.
        description: Optional description of the batch benchmark.
        tags: Optional list of tags for the batch benchmark.
        parallel: Whether to run benchmarks in parallel.
        max_workers: The maximum number of worker threads to use for parallel execution.
        
    Returns:
        BatchBenchmark: The created batch benchmark.
        
    Raises:
        ValueError: If the template ID is not found.
    """
    # Check if the template exists
    registry = BatchBenchmarkRegistry.get_instance()
    template = registry.get_template(template_id)
    
    if not template:
        raise ValueError(f"Template with ID {template_id} not found")
    
    # Generate a unique ID for the batch
    batch_id = f"batch_{uuid.uuid4().hex[:8]}"
    
    # Create the batch
    batch = BatchBenchmark(
        id=batch_id,
        name=name,
        model_ids=model_ids,
        template_id=template_id,
        description=description,
        tags=tags or [],
        parallel=parallel,
        max_workers=max_workers
    )
    
    # Register the batch
    registry.register_batch(batch)
    
    logger.info(f"Created batch benchmark: {batch_id}")
    return batch


def run_batch_benchmark(
    batch: BatchBenchmark,
    progress_callback: Optional[ProgressCallback] = None
) -> BatchResult:
    """
    Run a batch benchmark.
    
    Args:
        batch: The batch benchmark to run.
        progress_callback: Optional callback for tracking progress.
        
    Returns:
        BatchResult: The batch benchmark result.
    """
    # Create operation for tracking progress
    operation_id, progress_info = create_operation(
        OperationType.BENCHMARK, 
        operation_id=f"batch_{batch.id}",
        callback=progress_callback
    )
    
    try:
        # Start the operation
        start_operation(operation_id, f"Running batch benchmark: {batch.name}")
        
        # Get the template
        registry = BatchBenchmarkRegistry.get_instance()
        template = registry.get_template(batch.template_id)
        
        if not template:
            raise ValueError(f"Template with ID {batch.template_id} not found")
        
        # Update batch status
        batch.status = BatchStatus.RUNNING
        batch.started_at = datetime.datetime.now().isoformat()
        registry.update_batch(batch)
        
        # Initialize results
        benchmark_results: Dict[str, List[BenchmarkResult]] = {}
        
        # Get the total number of benchmarks to run
        total_benchmarks = len(batch.model_ids) * len(template.prompts)
        completed_benchmarks = 0
        
        # Record start time
        start_time = time.time()
        
        # Run benchmarks for each model
        if batch.parallel and batch.max_workers > 1:
            # Run benchmarks in parallel
            with concurrent.futures.ThreadPoolExecutor(max_workers=batch.max_workers) as executor:
                # Create a list of futures
                futures = []
                
                # Submit benchmarks for each model and prompt
                for model_id in batch.model_ids:
                    benchmark_results[model_id] = []
                    
                    for i, prompt in enumerate(template.prompts):
                        # Create benchmark
                        benchmark = ModelBenchmark(
                            model_id=model_id,
                            prompt=prompt,
                            max_tokens=template.max_tokens,
                            num_runs=template.num_runs,
                            temperature=template.temperature,
                            top_p=template.top_p,
                            tags=template.tags + batch.tags + [f"batch_{batch.id}"],
                            description=f"Batch benchmark: {batch.name} - Prompt {i+1}/{len(template.prompts)}",
                            save_logprobs=template.save_logprobs,
                            reference_text=template.reference_texts[i] if template.reference_texts and i < len(template.reference_texts) else None
                        )
                        
                        # Submit the benchmark
                        future = executor.submit(run_benchmark, benchmark)
                        futures.append((model_id, future))
                
                # Process futures as they complete
                for model_id, future in futures:
                    try:
                        result = future.result()
                        benchmark_results[model_id].append(result)
                        
                        # Update progress
                        completed_benchmarks += 1
                        progress = completed_benchmarks / total_benchmarks
                        update_progress(
                            operation_id,
                            progress,
                            f"Completed {completed_benchmarks}/{total_benchmarks} benchmarks"
                        )
                    except Exception as e:
                        logger.error(f"Error running benchmark for model {model_id}: {e}", exc_info=True)
                        # Continue with other benchmarks
        else:
            # Run benchmarks sequentially
            for model_id in batch.model_ids:
                benchmark_results[model_id] = []
                
                for i, prompt in enumerate(template.prompts):
                    # Update progress
                    progress = completed_benchmarks / total_benchmarks
                    update_progress(
                        operation_id,
                        progress,
                        f"Running benchmark for model {model_id} - Prompt {i+1}/{len(template.prompts)}"
                    )
                    
                    # Create benchmark
                    benchmark = ModelBenchmark(
                        model_id=model_id,
                        prompt=prompt,
                        max_tokens=template.max_tokens,
                        num_runs=template.num_runs,
                        temperature=template.temperature,
                        top_p=template.top_p,
                        tags=template.tags + batch.tags + [f"batch_{batch.id}"],
                        description=f"Batch benchmark: {batch.name} - Prompt {i+1}/{len(template.prompts)}",
                        save_logprobs=template.save_logprobs,
                        reference_text=template.reference_texts[i] if template.reference_texts and i < len(template.reference_texts) else None
                    )
                    
                    # Run benchmark
                    try:
                        result = run_benchmark(benchmark)
                        benchmark_results[model_id].append(result)
                    except Exception as e:
                        logger.error(f"Error running benchmark for model {model_id}: {e}", exc_info=True)
                        # Continue with other benchmarks
                    
                    # Update progress
                    completed_benchmarks += 1
        
        # Record end time
        end_time = time.time()
        duration_seconds = end_time - start_time
        
        # Create batch result
        batch_result = BatchResult(
            batch_id=batch.id,
            name=batch.name,
            template_id=batch.template_id,
            model_ids=batch.model_ids,
            benchmark_results=benchmark_results,
            started_at=batch.started_at,
            completed_at=datetime.datetime.now().isoformat(),
            duration_seconds=duration_seconds
        )
        
        # Update batch status
        batch.status = BatchStatus.COMPLETED
        batch.completed_at = batch_result.completed_at
        registry.update_batch(batch)
        
        # Register the result
        registry.register_result(batch_result)
        
        # Complete the operation
        complete_operation(
            operation_id,
            batch_result,
            f"Batch benchmark completed: {batch.name}"
        )
        
        return batch_result
    
    except Exception as e:
        logger.error(f"Error running batch benchmark: {e}", exc_info=True)
        
        # Update batch status
        batch.status = BatchStatus.FAILED
        batch.error = str(e)
        registry.update_batch(batch)
        
        # Create error result
        error_result = BatchResult(
            batch_id=batch.id,
            name=batch.name,
            template_id=batch.template_id,
            model_ids=batch.model_ids,
            benchmark_results={},
            started_at=batch.started_at or datetime.datetime.now().isoformat(),
            completed_at=datetime.datetime.now().isoformat(),
            duration_seconds=0.0,
            error=str(e)
        )
        
        # Register the error result
        registry.register_result(error_result)
        
        # Fail the operation
        fail_operation(operation_id, str(e), f"Failed to run batch benchmark: {batch.name}")
        
        return error_result


def get_batch_benchmarks() -> List[BatchBenchmark]:
    """
    Get all registered batch benchmarks.
    
    Returns:
        List[BatchBenchmark]: A list of batch benchmarks.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_batches()


def get_batch_results() -> List[BatchResult]:
    """
    Get all registered batch results.
    
    Returns:
        List[BatchResult]: A list of batch results.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_results()


def get_batch_result(batch_id: str) -> Optional[BatchResult]:
    """
    Get a batch result by ID.
    
    Args:
        batch_id: The ID of the batch to get the result for.
        
    Returns:
        Optional[BatchResult]: The batch result, or None if not found.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_result(batch_id)


def cancel_batch_benchmark(batch_id: str) -> bool:
    """
    Cancel a running batch benchmark.
    
    Args:
        batch_id: The ID of the batch to cancel.
        
    Returns:
        bool: True if the batch was cancelled, False otherwise.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    batch = registry.get_batch(batch_id)
    
    if not batch:
        logger.warning(f"Batch with ID {batch_id} not found")
        return False
    
    if batch.status != BatchStatus.RUNNING:
        logger.warning(f"Batch {batch_id} is not running (status: {batch.status.name})")
        return False
    
    # Update batch status
    batch.status = BatchStatus.CANCELLED
    batch.completed_at = datetime.datetime.now().isoformat()
    registry.update_batch(batch)
    
    logger.info(f"Cancelled batch benchmark: {batch_id}")
    return True


def export_batch_results(
    batch_result: BatchResult,
    format: str = "html",
    output_path: Optional[str] = None,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to a file.
    
    Args:
        batch_result: The batch result to export.
        format: The format to export to (html, pdf, pptx).
        output_path: Optional path to save the exported file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
        
    Raises:
        ValueError: If the format is not supported.
    """
    # Check if the format is supported
    if format not in ["html", "pdf", "pptx"]:
        raise ValueError(f"Unsupported export format: {format}")
    
    # Generate a default output path if not provided
    if not output_path:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"batch_benchmark_{batch_result.batch_id}_{timestamp}.{format}"
    
    # Create the export directory if it doesn't exist
    output_dir = os.path.dirname(output_path)
    if output_dir:
        ensure_directory(output_dir)
    
    # Export based on the format
    if format == "html":
        return _export_batch_results_html(batch_result, output_path, include_plots)
    elif format == "pdf":
        return _export_batch_results_pdf(batch_result, output_path, include_plots)
    elif format == "pptx":
        return _export_batch_results_pptx(batch_result, output_path, include_plots)
    
    # This should never happen due to the check above
    raise ValueError(f"Unsupported export format: {format}")


def _export_batch_results_html(
    batch_result: BatchResult,
    output_path: str,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to an HTML file.
    
    Args:
        batch_result: The batch result to export.
        output_path: The path to save the HTML file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
    """
    # Get the template
    registry = BatchBenchmarkRegistry.get_instance()
    template = registry.get_template(batch_result.template_id)
    
    if not template:
        logger.warning(f"Template with ID {batch_result.template_id} not found")
    
    # Check if visualization dependencies are available
    if include_plots and not check_visualization_dependencies():
        logger.warning("Visualization dependencies are not available, plots will not be included")
        include_plots = False
    
    # Create HTML content
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Batch Benchmark Results: {batch_result.name}</title>
        <style>
            body {{
                font-family: Arial, sans-serif;
                line-height: 1.6;
                color: #333;
                max-width: 1200px;
                margin: 0 auto;
                padding: 20px;
            }}
            h1, h2, h3, h4 {{
                color: #2c3e50;
            }}
            table {{
                border-collapse: collapse;
                width: 100%;
                margin-bottom: 20px;
            }}
            th, td {{
                border: 1px solid #ddd;
                padding: 8px;
                text-align: left;
            }}
            th {{
                background-color: #f2f2f2;
            }}
            tr:nth-child(even) {{
                background-color: #f9f9f9;
            }}
            .model-section {{
                margin-bottom: 30px;
                border: 1px solid #ddd;
                padding: 15px;
                border-radius: 5px;
            }}
            .prompt-section {{
                margin-bottom: 20px;
                padding: 10px;
                background-color: #f5f5f5;
                border-radius: 5px;
            }}
            .generated-text {{
                white-space: pre-wrap;
                background-color: #f8f8f8;
                padding: 10px;
                border-left: 3px solid #2c3e50;
                margin: 10px 0;
                font-family: monospace;
            }}
            .metrics-table {{
                width: 100%;
                margin-top: 10px;
            }}
            .plot-container {{
                margin: 20px 0;
                text-align: center;
            }}
            .plot-container img {{
                max-width: 100%;
                height: auto;
                border: 1px solid #ddd;
            }}
            .summary {{
                background-color: #e8f4f8;
                padding: 15px;
                border-radius: 5px;
                margin-bottom: 20px;
            }}
        </style>
    </head>
    <body>
        <h1>Batch Benchmark Results: {batch_result.name}</h1>
        
        <div class="summary">
            <h2>Summary</h2>
            <p><strong>Started:</strong> {batch_result.started_at}</p>
            <p><strong>Completed:</strong> {batch_result.completed_at}</p>
            <p><strong>Duration:</strong> {batch_result.duration_seconds:.2f} seconds</p>
            <p><strong>Models:</strong> {len(batch_result.model_ids)}</p>
            <p><strong>Prompts:</strong> {len(template.prompts) if template else "Unknown"}</p>
        </div>
    """
    
    # Add model comparison section if there are multiple models
    if len(batch_result.model_ids) > 1 and include_plots:
        html_content += """
        <h2>Model Comparison</h2>
        <div class="plot-container">
        """
        
        # Create comparison plots
        try:
            # Create a temporary directory for plots
            plots_dir = tempfile.mkdtemp()
            
            # Get all results for each model
            all_results = []
            for model_id, results in batch_result.benchmark_results.items():
                all_results.extend(results)
            
            # Create plots for different metrics
            metrics = ["avg_tokens_per_second", "avg_generation_time", "peak_memory_mb"]
            metric_names = ["Tokens per Second", "Generation Time (s)", "Memory Usage (MB)"]
            
            for metric, metric_name in zip(metrics, metric_names):
                # Create the plot
                plot_path = os.path.join(plots_dir, f"comparison_{metric}.png")
                
                # Sort ascending for metrics where lower is better
                sort_ascending = metric in ["avg_generation_time", "peak_memory_mb"]
                
                # Create the plot
                fig = plot_model_comparison(
                    all_results, 
                    metric, 
                    title=f"Model Comparison: {metric_name}", 
                    sort_ascending=sort_ascending
                )
                
                # Save the plot
                export_visualization(fig, plot_path)
                
                # Add the plot to the HTML
                plot_filename = os.path.basename(plot_path)
                html_content += f"""
                <h3>{metric_name}</h3>
                <img src="{plot_filename}" alt="Model comparison: {metric_name}">
                """
            
            # Create radar chart
            radar_path = os.path.join(plots_dir, "radar_chart.png")
            
            # Get model names
            model_names = {}
            for model_id in batch_result.model_ids:
                model_registry = ModelRegistry.get_instance()
                model_info = model_registry.get_model_info(model_id)
                model_names[model_id] = model_info.name if model_info else model_id
            
            # Create the radar chart
            fig = plot_model_radar_chart(
                all_results,
                list(model_names.values()),
                title="Model Comparison: Radar Chart"
            )
            
            # Save the plot
            export_visualization(fig, radar_path)
            
            # Add the plot to the HTML
            radar_filename = os.path.basename(radar_path)
            html_content += f"""
            <h3>Radar Chart</h3>
            <img src="{radar_filename}" alt="Model comparison: Radar Chart">
            """
            
            html_content += """
            </div>
            """
            
            # Copy plots to the output directory
            output_dir = os.path.dirname(output_path)
            for plot_file in os.listdir(plots_dir):
                plot_path = os.path.join(plots_dir, plot_file)
                import shutil
                shutil.copy2(plot_path, os.path.join(output_dir, plot_file))
        
        except Exception as e:
            logger.error(f"Error creating comparison plots: {e}", exc_info=True)
            html_content += f"""
            <p>Error creating comparison plots: {e}</p>
            </div>
            """
    
    # Add model sections
    for model_id in batch_result.model_ids:
        # Get model info
        model_registry = ModelRegistry.get_instance()
        model_info = model_registry.get_model_info(model_id)
        
        model_name = model_info.name if model_info else model_id
        model_type = model_info.model_type.name if model_info and hasattr(model_info, "model_type") else "Unknown"
        
        html_content += f"""
        <div class="model-section">
            <h2>Model: {model_name} ({model_id})</h2>
            <p><strong>Type:</strong> {model_type}</p>
        """
        
        # Add model results
        results = batch_result.benchmark_results.get(model_id, [])
        
        if not results:
            html_content += """
            <p>No results available for this model.</p>
            </div>
            """
            continue
        
        # Calculate average metrics
        avg_load_time = statistics.mean([r.load_time_seconds for r in results if r.load_time_seconds is not None])
        avg_generation_time = statistics.mean([r.avg_generation_time for r in results if r.avg_generation_time is not None])
        avg_tokens_per_second = statistics.mean([r.avg_tokens_per_second for r in results if r.avg_tokens_per_second is not None])
        avg_peak_memory = statistics.mean([r.peak_memory_mb for r in results if r.peak_memory_mb is not None])
        
        html_content += f"""
        <h3>Average Metrics</h3>
        <table class="metrics-table">
            <tr>
                <th>Load Time</th>
                <th>Generation Time</th>
                <th>Tokens per Second</th>
                <th>Memory Usage</th>
            </tr>
            <tr>
                <td>{avg_load_time:.2f} s</td>
                <td>{avg_generation_time:.2f} s</td>
                <td>{avg_tokens_per_second:.2f}</td>
                <td>{avg_peak_memory:.2f} MB</td>
            </tr>
        </table>
        """
        
        # Add prompt sections
        for i, result in enumerate(results):
            prompt = result.prompt
            
            html_content += f"""
            <div class="prompt-section">
                <h3>Prompt {i + 1}</h3>
                <p>{prompt}</p>
                
                <h4>Generated Text</h4>
            """
            
            if result.generated_texts:
                html_content += f"""
                <div class="generated-text">{result.generated_texts[0]}</div>
                """
            else:
                html_content += """
                <p>No generated text available.</p>
                """
            
            html_content += f"""
            <h4>Metrics</h4>
            <table class="metrics-table">
                <tr>
                    <th>Load Time</th>
                    <th>Generation Time</th>
                    <th>Tokens per Second</th>
                    <th>Memory Usage</th>
                </tr>
                <tr>
                    <td>{result.load_time_seconds:.2f} s</td>
                    <td>{result.avg_generation_time:.2f} s</td>
                    <td>{result.avg_tokens_per_second:.2f}</td>
                    <td>{result.peak_memory_mb:.2f} MB</td>
                </tr>
            </table>
            </div>
            """
        
        html_content += """
        </div>
        """
    
    # Close HTML
    html_content += """
    </body>
    </html>
    """
    
    # Write HTML to file
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    
    logger.info(f"Exported batch results to HTML: {output_path}")
    return output_path


def _export_batch_results_pdf(
    batch_result: BatchResult,
    output_path: str,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to a PDF file.
    
    Args:
        batch_result: The batch result to export.
        output_path: The path to save the PDF file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
    """
    try:
        # First export to HTML
        html_path = output_path.replace(".pdf", ".html")
        _export_batch_results_html(batch_result, html_path, include_plots)
        
        # Check if weasyprint is available
        try:
            from weasyprint import HTML
            
            # Convert HTML to PDF
            HTML(html_path).write_pdf(output_path)
            
            # Remove the temporary HTML file
            os.remove(html_path)
            
            logger.info(f"Exported batch results to PDF: {output_path}")
            return output_path
        
        except ImportError:
            logger.warning("weasyprint is not available, falling back to HTML export")
            return html_path
    
    except Exception as e:
        logger.error(f"Error exporting batch results to PDF: {e}", exc_info=True)
        raise


def _export_batch_results_pptx(
    batch_result: BatchResult,
    output_path: str,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to a PowerPoint file.
    
    Args:
        batch_result: The batch result to export.
        output_path: The path to save the PowerPoint file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
    """
    try:
        # Check if python-pptx is available
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Create a new presentation
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"Batch Benchmark Results: {batch_result.name}"
            subtitle.text = f"Generated on {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Add summary slide
            summary_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(summary_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Summary"
            
            # Get the template
            registry = BatchBenchmarkRegistry.get_instance()
            template = registry.get_template(batch_result.template_id)
            
            content.text = (
                f"Started: {batch_result.started_at}\n"
                f"Completed: {batch_result.completed_at}\n"
                f"Duration: {batch_result.duration_seconds:.2f} seconds\n"
                f"Models: {len(batch_result.model_ids)}\n"
                f"Prompts: {len(template.prompts) if template else 'Unknown'}"
            )
            
            # Add model comparison slides if there are multiple models and plots are included
            if len(batch_result.model_ids) > 1 and include_plots:
                # Check if visualization dependencies are available
                if check_visualization_dependencies():
                    # Create a temporary directory for plots
                    plots_dir = tempfile.mkdtemp()
                    
                    # Get all results for each model
                    all_results = []
                    for model_id, results in batch_result.benchmark_results.items():
                        all_results.extend(results)
                    
                    # Create plots for different metrics
                    metrics = ["avg_tokens_per_second", "avg_generation_time", "peak_memory_mb"]
                    metric_names = ["Tokens per Second", "Generation Time (s)", "Memory Usage (MB)"]
                    
                    for metric, metric_name in zip(metrics, metric_names):
                        # Create the plot
                        plot_path = os.path.join(plots_dir, f"comparison_{metric}.png")
                        
                        # Sort ascending for metrics where lower is better
                        sort_ascending = metric in ["avg_generation_time", "peak_memory_mb"]
                        
                        # Create the plot
                        fig = plot_model_comparison(
                            all_results, 
                            metric, 
                            title=f"Model Comparison: {metric_name}", 
                            sort_ascending=sort_ascending
                        )
                        
                        # Save the plot
                        export_visualization(fig, plot_path)
                        
                        # Add a slide for the plot
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        title = slide.shapes.title
                        title.text = f"Model Comparison: {metric_name}"
                        
                        # Add the image
                        slide.shapes.add_picture(plot_path, Inches(1), Inches(1.5), width=Inches(8))
                    
                    # Create radar chart
                    radar_path = os.path.join(plots_dir, "radar_chart.png")
                    
                    # Get model names
                    model_names = {}
                    for model_id in batch_result.model_ids:
                        model_registry = ModelRegistry.get_instance()
                        model_info = model_registry.get_model_info(model_id)
                        model_names[model_id] = model_info.name if model_info else model_id
                    
                    # Create the radar chart
                    fig = plot_model_radar_chart(
                        all_results,
                        list(model_names.values()),
                        title="Model Comparison: Radar Chart"
                    )
                    
                    # Save the plot
                    export_visualization(fig, radar_path)
                    
                    # Add a slide for the radar chart
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    title = slide.shapes.title
                    title.text = "Model Comparison: Radar Chart"
                    
                    # Add the image
                    slide.shapes.add_picture(radar_path, Inches(1), Inches(1.5), width=Inches(8))
            
            # Add model slides
            for model_id in batch_result.model_ids:
                # Get model info
                model_registry = ModelRegistry.get_instance()
                model_info = model_registry.get_model_info(model_id)
                
                model_name = model_info.name if model_info else model_id
                model_type = model_info.model_type.name if model_info and hasattr(model_info, "model_type") else "Unknown"
                
                # Add model slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Model: {model_name} ({model_id})"
                
                # Add model results
                results = batch_result.benchmark_results.get(model_id, [])
                
                if not results:
                    content.text = "No results available for this model."
                    continue
                
                # Calculate average metrics
                avg_load_time = statistics.mean([r.load_time_seconds for r in results if r.load_time_seconds is not None])
                avg_generation_time = statistics.mean([r.avg_generation_time for r in results if r.avg_generation_time is not None])
                avg_tokens_per_second = statistics.mean([r.avg_tokens_per_second for r in results if r.avg_tokens_per_second is not None])
                avg_peak_memory = statistics.mean([r.peak_memory_mb for r in results if r.peak_memory_mb is not None])
                
                content.text = (
                    f"Type: {model_type}\n\n"
                    f"Average Metrics:\n"
                    f"Load Time: {avg_load_time:.2f} s\n"
                    f"Generation Time: {avg_generation_time:.2f} s\n"
                    f"Tokens per Second: {avg_tokens_per_second:.2f}\n"
                    f"Memory Usage: {avg_peak_memory:.2f} MB"
                )
                
                # Add prompt slides
                for i, result in enumerate(results):
                    prompt = result.prompt
                    
                    # Add prompt slide
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = f"{model_name}: Prompt {i + 1}"
                    
                    # Add prompt and metrics
                    content_text = (
                        f"Prompt: {prompt}\n\n"
                        f"Metrics:\n"
                        f"Load Time: {result.load_time_seconds:.2f} s\n"
                        f"Generation Time: {result.avg_generation_time:.2f} s\n"
                        f"Tokens per Second: {result.avg_tokens_per_second:.2f}\n"
                        f"Memory Usage: {result.peak_memory_mb:.2f} MB\n\n"
                    )
                    
                    # Add generated text
                    if result.generated_texts:
                        # Truncate if too long
                        generated_text = result.generated_texts[0]
                        if len(generated_text) > 500:
                            generated_text = generated_text[:500] + "..."
                        
                        content_text += f"Generated Text:\n{generated_text}"
                    else:
                        content_text += "No generated text available."
                    
                    content.text = content_text
            
            # Save the presentation
            prs.save(output_path)
            
            logger.info(f"Exported batch results to PowerPoint: {output_path}")
            return output_path
        
        except ImportError:
            logger.warning("python-pptx is not available, falling back to HTML export")
            html_path = output_path.replace(".pptx", ".html")
            return _export_batch_results_html(batch_result, html_path, include_plots)
    
    except Exception as e:
        logger.error(f"Error exporting batch results to PowerPoint: {e}", exc_info=True)
        raise


def create_predefined_benchmark_templates() -> List[str]:
    """
    Create a set of predefined benchmark templates for common use cases.
    
    Returns:
        List[str]: A list of template IDs.
    """
    # Define predefined templates
    predefined_templates = [
        {
            "name": "Creative Writing",
            "prompts": [
                "Write a short story about a robot learning to paint.",
                "Write a poem about the changing seasons.",
                "Write a dialogue between two characters discussing the ethics of AI."
            ],
            "max_tokens": 150,
            "tags": ["creative", "writing"]
        },
        {
            "name": "Question Answering",
            "prompts": [
                "What are the main challenges in developing sustainable energy sources?",
                "Explain the concept of quantum computing to a high school student.",
                "What are the ethical implications of artificial intelligence in healthcare?"
            ],
            "max_tokens": 200,
            "tags": ["qa", "factual"]
        },
        {
            "name": "Summarization",
            "prompts": [
                "Summarize the following text: The development of artificial intelligence has accelerated in recent years, with breakthroughs in natural language processing, computer vision, and reinforcement learning. These advances have led to applications in healthcare, finance, transportation, and many other fields. However, they also raise important questions about privacy, bias, job displacement, and the long-term implications of increasingly autonomous systems.",
                "Summarize the following text: Climate change is one of the most pressing challenges facing humanity today. Rising global temperatures are leading to more frequent and severe weather events, sea level rise, and disruptions to ecosystems worldwide. Addressing this challenge requires a combination of mitigation strategies to reduce greenhouse gas emissions and adaptation measures to cope with the changes that are already underway.",
                "Summarize the following text: The human genome project was an international scientific research project with the goal of determining the base pairs that make up human DNA, and of identifying and mapping all of the genes of the human genome. It remains the world's largest collaborative biological project. The project was completed in 2003, 13 years after its launch, and made the sequence of the human genome available to researchers worldwide."
            ],
            "max_tokens": 100,
            "tags": ["summarization", "compression"]
        },
        {
            "name": "Code Generation",
            "prompts": [
                "Write a Python function that calculates the Fibonacci sequence up to n terms.",
                "Write a JavaScript function to sort an array of objects by a specific property.",
                "Write a SQL query to find the top 5 customers who have placed the most orders."
            ],
            "max_tokens": 200,
            "tags": ["code", "programming"]
        },
        {
            "name": "Reasoning",
            "prompts": [
                "A bat and a ball cost $1.10 in total. The bat costs $1.00 more than the ball. How much does the ball cost?",
                "If it takes 5 machines 5 minutes to make 5 widgets, how long would it take 100 machines to make 100 widgets?",
                "In a lake, there is a patch of lily pads. Every day, the patch doubles in size. If it takes 48 days for the patch to cover the entire lake, how long would it take for the patch to cover half of the lake?"
            ],
            "max_tokens": 150,
            "tags": ["reasoning", "logic"]
        }
    ]
    
    # Create templates
    template_ids = []
    
    for template_data in predefined_templates:
        template = create_benchmark_template(
            name=template_data["name"],
            prompts=template_data["prompts"],
            max_tokens=template_data["max_tokens"],
            tags=template_data["tags"],
            description=f"Predefined template for {template_data['name'].lower()} tasks"
        )
        
        template_ids.append(template.id)
    
    logger.info(f"Created {len(template_ids)} predefined benchmark templates")
    return template_ids


if __name__ == "__main__":
    # Example usage
    import argparse
    
    parser = argparse.ArgumentParser(description="Batch Benchmarking Tool")
    parser.add_argument("--create-templates", action="store_true", help="Create predefined benchmark templates")
    parser.add_argument("--list-templates", action="store_true", help="List available benchmark templates")
    parser.add_argument("--create-batch", action="store_true", help="Create a batch benchmark")
    parser.add_argument("--template", type=str, help="Template ID to use for the batch benchmark")
    parser.add_argument("--models", type=str, nargs="+", help="Model IDs to benchmark")
    parser.add_argument("--name", type=str, help="Name for the batch benchmark")
    parser.add_argument("--run", type=str, help="Batch ID to run")
    parser.add_argument("--list-batches", action="store_true", help="List available batch benchmarks")
    parser.add_argument("--export", type=str, help="Batch ID to export")
    parser.add_argument("--format", type=str, choices=["html", "pdf", "pptx"], default="html", help="Export format")
    parser.add_argument("--output", type=str, help="Output path for export")
    
    args = parser.parse_args()
    
    if args.create_templates:
        print("Creating predefined benchmark templates...")
        template_ids = create_predefined_benchmark_templates()
        print(f"Created {len(template_ids)} predefined benchmark templates")
    
    elif args.list_templates:
        print("Available benchmark templates:")
        templates = get_benchmark_templates()
        
        for template in templates:
            print(f"\nTemplate: {template.name} ({template.id})")
            print(f"Prompts: {len(template.prompts)}")
            print(f"Max Tokens: {template.max_tokens}")
            print(f"Tags: {', '.join(template.tags)}")
            if template.description:
                print(f"Description: {template.description}")
    
    elif args.create_batch:
        if not args.template:
            print("Error: Template ID is required")
            sys.exit(1)
        
        if not args.models:
            print("Error: At least one model ID is required")
            sys.exit(1)
        
        if not args.name:
            print("Error: Batch name is required")
            sys.exit(1)
        
        print(f"Creating batch benchmark: {args.name}")
        batch = create_batch_benchmark(
            name=args.name,
            model_ids=args.models,
            template_id=args.template
        )
        
        print(f"Created batch benchmark: {batch.id}")
    
    elif args.run:
        print(f"Running batch benchmark: {args.run}")
        
        registry = BatchBenchmarkRegistry.get_instance()
        batch = registry.get_batch(args.run)
        
        if not batch:
            print(f"Error: Batch with ID {args.run} not found")
            sys.exit(1)
        
        result = run_batch_benchmark(batch)
        
        print(f"Batch benchmark completed: {result.batch_id}")
        print(f"Duration: {result.duration_seconds:.2f} seconds")
        
        for model_id, results in result.benchmark_results.items():
            print(f"\nModel: {model_id}")
            print(f"Results: {len(results)}")
    
    elif args.list_batches:
        print("Available batch benchmarks:")
        batches = get_batch_benchmarks()
        
        for batch in batches:
            print(f"\nBatch: {batch.name} ({batch.id})")
            print(f"Models: {len(batch.model_ids)}")
            print(f"Template: {batch.template_id}")
            print(f"Status: {batch.status.name}")
            if batch.description:
                print(f"Description: {batch.description}")
    
    elif args.export:
        print(f"Exporting batch results: {args.export}")
        
        result = get_batch_result(args.export)
        
        if not result:
            print(f"Error: Batch result with ID {args.export} not found")
            sys.exit(1)
        
        output_path = args.output or f"batch_results_{args.export}.{args.format}"
        
        export_path = export_batch_results(
            result,
            format=args.format,
            output_path=output_path
        )
        
        print(f"Exported batch results to: {export_path}")
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Batch Benchmarking for RebelSCRIBE.

This module provides functionality for running multiple benchmarks in batch mode,
managing benchmark templates, and generating comprehensive reports.

Example usage:
    ```python
    from src.ai.batch_benchmarking import (
        BatchBenchmark, BenchmarkTemplate, run_batch_benchmark,
        create_benchmark_template, get_benchmark_templates,
        export_batch_results
    )
    
    # Create a benchmark template
    template = create_benchmark_template(
        name="Creative Writing",
        prompts=[
            "Write a short story about a robot learning to paint.",
            "Write a poem about the changing seasons.",
            "Write a dialogue between two characters discussing the ethics of AI."
        ],
        max_tokens=150,
        num_runs=3,
        tags=["creative", "writing"]
    )
    
    # Create a batch benchmark
    batch = BatchBenchmark(
        name="Creative Writing Benchmark",
        model_ids=["llama-2-7b", "mistral-7b", "phi-2"],
        template_id=template.id,
        description="Comparing models on creative writing tasks"
    )
    
    # Run the batch benchmark
    results = run_batch_benchmark(batch)
    
    # Export the results
    export_batch_results(results, format="pdf", output_path="benchmark_results.pdf")
    ```
"""

import os
import json
import time
import logging
import threading
import statistics
import uuid
import tempfile
from enum import Enum, auto
from typing import Dict, List, Optional, Any, Set, Tuple, Union, Callable
from dataclasses import dataclass, field, asdict
import datetime
import concurrent.futures
from pathlib import Path

from src.utils.logging_utils import get_logger
from src.utils.file_utils import ensure_directory
from src.ai.progress_callbacks import (
    OperationType, ProgressStatus, ProgressInfo, ProgressCallback,
    create_operation, start_operation, update_progress,
    complete_operation, fail_operation
)
from src.ai.model_benchmarking import (
    ModelBenchmark, BenchmarkResult, run_benchmark,
    get_benchmark_results, BenchmarkRegistry
)
from src.ai.model_registry import (
    ModelRegistry, ModelInfo, ModelType, ModelSource
)
from src.ai.benchmark_visualization import (
    plot_model_comparison, plot_benchmark_history, plot_metric_correlation,
    plot_model_radar_chart, export_visualization, generate_benchmark_report_html,
    check_visualization_dependencies
)

logger = get_logger(__name__)


class BatchStatus(Enum):
    """Enum representing the status of a batch benchmark."""
    PENDING = auto()
    RUNNING = auto()
    COMPLETED = auto()
    FAILED = auto()
    CANCELLED = auto()


@dataclass
class BenchmarkTemplate:
    """Class for defining a benchmark template."""
    id: str
    name: str
    prompts: List[str]
    max_tokens: int = 100
    num_runs: int = 3
    temperature: float = 0.7
    top_p: float = 0.9
    tags: List[str] = field(default_factory=list)
    description: Optional[str] = None
    reference_texts: Optional[List[str]] = None
    save_logprobs: bool = True
    created_at: str = field(default_factory=lambda: datetime.datetime.now().isoformat())
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the template to a dictionary."""
        return asdict(self)
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BenchmarkTemplate':
        """Create a BenchmarkTemplate instance from a dictionary."""
        return cls(**data)


@dataclass
class BatchBenchmark:
    """Class for defining a batch benchmark."""
    id: str
    name: str
    model_ids: List[str]
    template_id: str
    description: Optional[str] = None
    tags: List[str] = field(default_factory=list)
    parallel: bool = False
    max_workers: int = 1
    status: BatchStatus = BatchStatus.PENDING
    created_at: str = field(default_factory=lambda: datetime.datetime.now().isoformat())
    started_at: Optional[str] = None
    completed_at: Optional[str] = None
    error: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the batch benchmark to a dictionary."""
        data = asdict(self)
        # Convert enum to string
        data["status"] = self.status.name
        return data
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BatchBenchmark':
        """Create a BatchBenchmark instance from a dictionary."""
        # Convert string status back to enum
        status_str = data.pop("status", "PENDING")
        status = BatchStatus[status_str]
        
        return cls(
            status=status,
            **data
        )


@dataclass
class BatchResult:
    """Class for storing batch benchmark results."""
    batch_id: str
    name: str
    template_id: str
    model_ids: List[str]
    benchmark_results: Dict[str, List[BenchmarkResult]]
    started_at: str
    completed_at: str
    duration_seconds: float
    error: Optional[str] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert the batch result to a dictionary."""
        data = asdict(self)
        # Convert BenchmarkResult objects to dictionaries
        data["benchmark_results"] = {
            model_id: [result.to_dict() for result in results]
            for model_id, results in self.benchmark_results.items()
        }
        return data
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'BatchResult':
        """Create a BatchResult instance from a dictionary."""
        # Convert dictionaries back to BenchmarkResult objects
        benchmark_results_dict = data.pop("benchmark_results", {})
        benchmark_results = {
            model_id: [BenchmarkResult.from_dict(result_dict) for result_dict in results_dict]
            for model_id, results_dict in benchmark_results_dict.items()
        }
        
        return cls(
            benchmark_results=benchmark_results,
            **data
        )


class BatchBenchmarkRegistry:
    """
    Singleton class for managing batch benchmarks, templates, and results.
    
    This class provides functionality for storing, retrieving, and managing
    batch benchmark definitions, templates, and results.
    """
    
    _instance = None
    _lock = threading.Lock()
    
    @classmethod
    def get_instance(cls) -> 'BatchBenchmarkRegistry':
        """Get the singleton instance of the BatchBenchmarkRegistry."""
        with cls._lock:
            if cls._instance is None:
                cls._instance = cls()
            return cls._instance
    
    def __init__(self):
        """Initialize the batch benchmark registry."""
        # Ensure this is only called once
        if BatchBenchmarkRegistry._instance is not None:
            raise RuntimeError("Use BatchBenchmarkRegistry.get_instance() to get the singleton instance")
        
        # Initialize registry data
        self._templates: Dict[str, BenchmarkTemplate] = {}
        self._batches: Dict[str, BatchBenchmark] = {}
        self._results: Dict[str, BatchResult] = {}
        self._registry_file = self._get_registry_file_path()
        
        # Load existing registry data
        self._load_registry()
    
    def _get_registry_file_path(self) -> str:
        """Get the path to the registry file."""
        # Get the registry directory
        registry_dir = os.environ.get(
            "REBELSCRIBE_BENCHMARK_DIR", 
            os.path.join(os.path.expanduser("~"), ".rebelscribe", "benchmarks")
        )
        
        # Ensure the directory exists
        ensure_directory(registry_dir)
        
        # Return the path to the registry file
        return os.path.join(registry_dir, "batch_benchmark_registry.json")
    
    def _load_registry(self) -> None:
        """Load the registry data from the registry file."""
        if not os.path.exists(self._registry_file):
            logger.info(f"Registry file not found at {self._registry_file}, creating new registry")
            return
        
        try:
            with open(self._registry_file, "r") as f:
                data = json.load(f)
            
            # Convert dictionaries to BenchmarkTemplate objects
            for template_id, template_data in data.get("templates", {}).items():
                self._templates[template_id] = BenchmarkTemplate.from_dict(template_data)
            
            # Convert dictionaries to BatchBenchmark objects
            for batch_id, batch_data in data.get("batches", {}).items():
                self._batches[batch_id] = BatchBenchmark.from_dict(batch_data)
            
            # Convert dictionaries to BatchResult objects
            for result_id, result_data in data.get("results", {}).items():
                self._results[result_id] = BatchResult.from_dict(result_data)
            
            logger.info(
                f"Loaded {len(self._templates)} templates, {len(self._batches)} batches, "
                f"and {len(self._results)} results from registry"
            )
        
        except Exception as e:
            logger.error(f"Error loading registry: {e}", exc_info=True)
            # Create a backup of the corrupted file
            if os.path.exists(self._registry_file):
                backup_path = f"{self._registry_file}.bak.{int(time.time())}"
                try:
                    import shutil
                    shutil.copy2(self._registry_file, backup_path)
                    logger.info(f"Created backup of corrupted registry at {backup_path}")
                except Exception as backup_error:
                    logger.error(f"Error creating backup: {backup_error}")
    
    def _save_registry(self) -> None:
        """Save the registry data to the registry file."""
        try:
            # Convert objects to dictionaries
            templates_dict = {template_id: template.to_dict() for template_id, template in self._templates.items()}
            batches_dict = {batch_id: batch.to_dict() for batch_id, batch in self._batches.items()}
            results_dict = {result_id: result.to_dict() for result_id, result in self._results.items()}
            
            # Create the data to save
            data = {
                "templates": templates_dict,
                "batches": batches_dict,
                "results": results_dict,
                "last_updated": datetime.datetime.now().isoformat()
            }
            
            # Save to file
            with open(self._registry_file, "w") as f:
                json.dump(data, f, indent=2)
            
            logger.info(
                f"Saved {len(self._templates)} templates, {len(self._batches)} batches, "
                f"and {len(self._results)} results to registry"
            )
        
        except Exception as e:
            logger.error(f"Error saving registry: {e}", exc_info=True)
    
    def register_template(self, template: BenchmarkTemplate) -> str:
        """
        Register a benchmark template in the registry.
        
        Args:
            template: The benchmark template to register.
            
        Returns:
            str: The ID of the registered template.
        """
        with self._lock:
            # Use the provided ID or generate a new one
            template_id = template.id
            
            # Register the template
            self._templates[template_id] = template
            self._save_registry()
            
            logger.info(f"Registered benchmark template: {template_id}")
            return template_id
    
    def get_template(self, template_id: str) -> Optional[BenchmarkTemplate]:
        """
        Get a benchmark template by ID.
        
        Args:
            template_id: The ID of the template to get.
            
        Returns:
            Optional[BenchmarkTemplate]: The template, or None if not found.
        """
        return self._templates.get(template_id)
    
    def get_all_templates(self) -> List[BenchmarkTemplate]:
        """
        Get all registered benchmark templates.
        
        Returns:
            List[BenchmarkTemplate]: A list of benchmark templates.
        """
        return list(self._templates.values())
    
    def register_batch(self, batch: BatchBenchmark) -> str:
        """
        Register a batch benchmark in the registry.
        
        Args:
            batch: The batch benchmark to register.
            
        Returns:
            str: The ID of the registered batch.
        """
        with self._lock:
            # Use the provided ID or generate a new one
            batch_id = batch.id
            
            # Register the batch
            self._batches[batch_id] = batch
            self._save_registry()
            
            logger.info(f"Registered batch benchmark: {batch_id}")
            return batch_id
    
    def update_batch(self, batch: BatchBenchmark) -> None:
        """
        Update a batch benchmark in the registry.
        
        Args:
            batch: The batch benchmark to update.
        """
        with self._lock:
            # Update the batch
            self._batches[batch.id] = batch
            self._save_registry()
            
            logger.info(f"Updated batch benchmark: {batch.id}")
    
    def get_batch(self, batch_id: str) -> Optional[BatchBenchmark]:
        """
        Get a batch benchmark by ID.
        
        Args:
            batch_id: The ID of the batch to get.
            
        Returns:
            Optional[BatchBenchmark]: The batch, or None if not found.
        """
        return self._batches.get(batch_id)
    
    def get_all_batches(self) -> List[BatchBenchmark]:
        """
        Get all registered batch benchmarks.
        
        Returns:
            List[BatchBenchmark]: A list of batch benchmarks.
        """
        return list(self._batches.values())
    
    def register_result(self, result: BatchResult) -> str:
        """
        Register a batch result in the registry.
        
        Args:
            result: The batch result to register.
            
        Returns:
            str: The ID of the registered result.
        """
        with self._lock:
            # Use the batch ID as the result ID
            result_id = result.batch_id
            
            # Register the result
            self._results[result_id] = result
            self._save_registry()
            
            logger.info(f"Registered batch result: {result_id}")
            return result_id
    
    def get_result(self, result_id: str) -> Optional[BatchResult]:
        """
        Get a batch result by ID.
        
        Args:
            result_id: The ID of the result to get.
            
        Returns:
            Optional[BatchResult]: The result, or None if not found.
        """
        return self._results.get(result_id)
    
    def get_all_results(self) -> List[BatchResult]:
        """
        Get all registered batch results.
        
        Returns:
            List[BatchResult]: A list of batch results.
        """
        return list(self._results.values())


def create_benchmark_template(
    name: str,
    prompts: List[str],
    max_tokens: int = 100,
    num_runs: int = 3,
    temperature: float = 0.7,
    top_p: float = 0.9,
    tags: Optional[List[str]] = None,
    description: Optional[str] = None,
    reference_texts: Optional[List[str]] = None,
    save_logprobs: bool = True
) -> BenchmarkTemplate:
    """
    Create a benchmark template.
    
    Args:
        name: The name of the template.
        prompts: The list of prompts to use for benchmarking.
        max_tokens: The maximum number of tokens to generate.
        num_runs: The number of times to run each benchmark.
        temperature: The temperature to use for generation.
        top_p: The top-p value to use for generation.
        tags: Optional list of tags for the template.
        description: Optional description of the template.
        reference_texts: Optional list of reference texts for BLEU score calculation.
        save_logprobs: Whether to save token logprobs for perplexity calculation.
        
    Returns:
        BenchmarkTemplate: The created benchmark template.
    """
    # Generate a unique ID for the template
    template_id = f"template_{uuid.uuid4().hex[:8]}"
    
    # Create the template
    template = BenchmarkTemplate(
        id=template_id,
        name=name,
        prompts=prompts,
        max_tokens=max_tokens,
        num_runs=num_runs,
        temperature=temperature,
        top_p=top_p,
        tags=tags or [],
        description=description,
        reference_texts=reference_texts,
        save_logprobs=save_logprobs
    )
    
    # Register the template
    registry = BatchBenchmarkRegistry.get_instance()
    registry.register_template(template)
    
    logger.info(f"Created benchmark template: {template_id}")
    return template


def get_benchmark_templates() -> List[BenchmarkTemplate]:
    """
    Get all registered benchmark templates.
    
    Returns:
        List[BenchmarkTemplate]: A list of benchmark templates.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_templates()


def create_batch_benchmark(
    name: str,
    model_ids: List[str],
    template_id: str,
    description: Optional[str] = None,
    tags: Optional[List[str]] = None,
    parallel: bool = False,
    max_workers: int = 1
) -> BatchBenchmark:
    """
    Create a batch benchmark.
    
    Args:
        name: The name of the batch benchmark.
        model_ids: The list of model IDs to benchmark.
        template_id: The ID of the template to use.
        description: Optional description of the batch benchmark.
        tags: Optional list of tags for the batch benchmark.
        parallel: Whether to run benchmarks in parallel.
        max_workers: The maximum number of worker threads to use for parallel execution.
        
    Returns:
        BatchBenchmark: The created batch benchmark.
        
    Raises:
        ValueError: If the template ID is not found.
    """
    # Check if the template exists
    registry = BatchBenchmarkRegistry.get_instance()
    template = registry.get_template(template_id)
    
    if not template:
        raise ValueError(f"Template with ID {template_id} not found")
    
    # Generate a unique ID for the batch
    batch_id = f"batch_{uuid.uuid4().hex[:8]}"
    
    # Create the batch
    batch = BatchBenchmark(
        id=batch_id,
        name=name,
        model_ids=model_ids,
        template_id=template_id,
        description=description,
        tags=tags or [],
        parallel=parallel,
        max_workers=max_workers
    )
    
    # Register the batch
    registry.register_batch(batch)
    
    logger.info(f"Created batch benchmark: {batch_id}")
    return batch


def run_batch_benchmark(
    batch: BatchBenchmark,
    progress_callback: Optional[ProgressCallback] = None
) -> BatchResult:
    """
    Run a batch benchmark.
    
    Args:
        batch: The batch benchmark to run.
        progress_callback: Optional callback for tracking progress.
        
    Returns:
        BatchResult: The batch benchmark result.
    """
    # Create operation for tracking progress
    operation_id, progress_info = create_operation(
        OperationType.BENCHMARK, 
        operation_id=f"batch_{batch.id}",
        callback=progress_callback
    )
    
    try:
        # Start the operation
        start_operation(operation_id, f"Running batch benchmark: {batch.name}")
        
        # Get the template
        registry = BatchBenchmarkRegistry.get_instance()
        template = registry.get_template(batch.template_id)
        
        if not template:
            raise ValueError(f"Template with ID {batch.template_id} not found")
        
        # Update batch status
        batch.status = BatchStatus.RUNNING
        batch.started_at = datetime.datetime.now().isoformat()
        registry.update_batch(batch)
        
        # Initialize results
        benchmark_results: Dict[str, List[BenchmarkResult]] = {}
        
        # Get the total number of benchmarks to run
        total_benchmarks = len(batch.model_ids) * len(template.prompts)
        completed_benchmarks = 0
        
        # Record start time
        start_time = time.time()
        
        # Run benchmarks for each model
        if batch.parallel and batch.max_workers > 1:
            # Run benchmarks in parallel
            with concurrent.futures.ThreadPoolExecutor(max_workers=batch.max_workers) as executor:
                # Create a list of futures
                futures = []
                
                # Submit benchmarks for each model and prompt
                for model_id in batch.model_ids:
                    benchmark_results[model_id] = []
                    
                    for i, prompt in enumerate(template.prompts):
                        # Create benchmark
                        benchmark = ModelBenchmark(
                            model_id=model_id,
                            prompt=prompt,
                            max_tokens=template.max_tokens,
                            num_runs=template.num_runs,
                            temperature=template.temperature,
                            top_p=template.top_p,
                            tags=template.tags + batch.tags + [f"batch_{batch.id}"],
                            description=f"Batch benchmark: {batch.name} - Prompt {i+1}/{len(template.prompts)}",
                            save_logprobs=template.save_logprobs,
                            reference_text=template.reference_texts[i] if template.reference_texts and i < len(template.reference_texts) else None
                        )
                        
                        # Submit the benchmark
                        future = executor.submit(run_benchmark, benchmark)
                        futures.append((model_id, future))
                
                # Process futures as they complete
                for model_id, future in futures:
                    try:
                        result = future.result()
                        benchmark_results[model_id].append(result)
                        
                        # Update progress
                        completed_benchmarks += 1
                        progress = completed_benchmarks / total_benchmarks
                        update_progress(
                            operation_id,
                            progress,
                            f"Completed {completed_benchmarks}/{total_benchmarks} benchmarks"
                        )
                    except Exception as e:
                        logger.error(f"Error running benchmark for model {model_id}: {e}", exc_info=True)
                        # Continue with other benchmarks
        else:
            # Run benchmarks sequentially
            for model_id in batch.model_ids:
                benchmark_results[model_id] = []
                
                for i, prompt in enumerate(template.prompts):
                    # Update progress
                    progress = completed_benchmarks / total_benchmarks
                    update_progress(
                        operation_id,
                        progress,
                        f"Running benchmark for model {model_id} - Prompt {i+1}/{len(template.prompts)}"
                    )
                    
                    # Create benchmark
                    benchmark = ModelBenchmark(
                        model_id=model_id,
                        prompt=prompt,
                        max_tokens=template.max_tokens,
                        num_runs=template.num_runs,
                        temperature=template.temperature,
                        top_p=template.top_p,
                        tags=template.tags + batch.tags + [f"batch_{batch.id}"],
                        description=f"Batch benchmark: {batch.name} - Prompt {i+1}/{len(template.prompts)}",
                        save_logprobs=template.save_logprobs,
                        reference_text=template.reference_texts[i] if template.reference_texts and i < len(template.reference_texts) else None
                    )
                    
                    # Run benchmark
                    try:
                        result = run_benchmark(benchmark)
                        benchmark_results[model_id].append(result)
                    except Exception as e:
                        logger.error(f"Error running benchmark for model {model_id}: {e}", exc_info=True)
                        # Continue with other benchmarks
                    
                    # Update progress
                    completed_benchmarks += 1
        
        # Record end time
        end_time = time.time()
        duration_seconds = end_time - start_time
        
        # Create batch result
        batch_result = BatchResult(
            batch_id=batch.id,
            name=batch.name,
            template_id=batch.template_id,
            model_ids=batch.model_ids,
            benchmark_results=benchmark_results,
            started_at=batch.started_at,
            completed_at=datetime.datetime.now().isoformat(),
            duration_seconds=duration_seconds
        )
        
        # Update batch status
        batch.status = BatchStatus.COMPLETED
        batch.completed_at = batch_result.completed_at
        registry.update_batch(batch)
        
        # Register the result
        registry.register_result(batch_result)
        
        # Complete the operation
        complete_operation(
            operation_id,
            batch_result,
            f"Batch benchmark completed: {batch.name}"
        )
        
        return batch_result
    
    except Exception as e:
        logger.error(f"Error running batch benchmark: {e}", exc_info=True)
        
        # Update batch status
        batch.status = BatchStatus.FAILED
        batch.error = str(e)
        registry.update_batch(batch)
        
        # Create error result
        error_result = BatchResult(
            batch_id=batch.id,
            name=batch.name,
            template_id=batch.template_id,
            model_ids=batch.model_ids,
            benchmark_results={},
            started_at=batch.started_at or datetime.datetime.now().isoformat(),
            completed_at=datetime.datetime.now().isoformat(),
            duration_seconds=0.0,
            error=str(e)
        )
        
        # Register the error result
        registry.register_result(error_result)
        
        # Fail the operation
        fail_operation(operation_id, str(e), f"Failed to run batch benchmark: {batch.name}")
        
        return error_result


def get_batch_benchmarks() -> List[BatchBenchmark]:
    """
    Get all registered batch benchmarks.
    
    Returns:
        List[BatchBenchmark]: A list of batch benchmarks.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_batches()


def get_batch_results() -> List[BatchResult]:
    """
    Get all registered batch results.
    
    Returns:
        List[BatchResult]: A list of batch results.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_all_results()


def get_batch_result(batch_id: str) -> Optional[BatchResult]:
    """
    Get a batch result by ID.
    
    Args:
        batch_id: The ID of the batch to get the result for.
        
    Returns:
        Optional[BatchResult]: The batch result, or None if not found.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    return registry.get_result(batch_id)


def cancel_batch_benchmark(batch_id: str) -> bool:
    """
    Cancel a running batch benchmark.
    
    Args:
        batch_id: The ID of the batch to cancel.
        
    Returns:
        bool: True if the batch was cancelled, False otherwise.
    """
    registry = BatchBenchmarkRegistry.get_instance()
    batch = registry.get_batch(batch_id)
    
    if not batch:
        logger.warning(f"Batch with ID {batch_id} not found")
        return False
    
    if batch.status != BatchStatus.RUNNING:
        logger.warning(f"Batch {batch_id} is not running (status: {batch.status.name})")
        return False
    
    # Update batch status
    batch.status = BatchStatus.CANCELLED
    batch.completed_at = datetime.datetime.now().isoformat()
    registry.update_batch(batch)
    
    logger.info(f"Cancelled batch benchmark: {batch_id}")
    return True


def export_batch_results(
    batch_result: BatchResult,
    format: str = "html",
    output_path: Optional[str] = None,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to a file.
    
    Args:
        batch_result: The batch result to export.
        format: The format to export to (html, pdf, pptx).
        output_path: Optional path to save the exported file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
        
    Raises:
        ValueError: If the format is not supported.
    """
    # Check if the format is supported
    if format not in ["html", "pdf", "pptx"]:
        raise ValueError(f"Unsupported export format: {format}")
    
    # Generate a default output path if not provided
    if not output_path:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"batch_benchmark_{batch_result.batch_id}_{timestamp}.{format}"
    
    # Create the export directory if it doesn't exist
    output_dir = os.path.dirname(output_path)
    if output_dir:
        ensure_directory(output_dir)
    
    # Export based on the format
    if format == "html":
        return _export_batch_results_html(batch_result, output_path, include_plots)
    elif format == "pdf":
        return _export_batch_results_pdf(batch_result, output_path, include_plots)
    elif format == "pptx":
        return _export_batch_results_pptx(batch_result, output_path, include_plots)
    
    # This should never happen due to the check above
    raise ValueError(f"Unsupported export format: {format}")


def _export_batch_results_html(
    batch_result: BatchResult,
    output_path: str,
    include_plots: bool = True
) -> str:
    """
    Export batch benchmark results to an HTML file.
    
    Args:
        batch_result: The batch result to export.
        output_path: The path to save the HTML file.
        include_plots: Whether to include plots in the export.
        
    Returns:
        str: The path to the exported file.
    """
    # Get the template
    registry = BatchBenchmarkRegistry.get_instance()
